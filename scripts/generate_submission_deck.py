from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "submission"
REPO_URL = "https://github.com/HyeokyYun/26ACEProgram_TechDigitalOmni.git"

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(246, 246, 246)
MID = RGBColor(100, 100, 100)
GREEN = RGBColor(38, 122, 72)
RED = RGBColor(180, 40, 40)
BLUE = RGBColor(20, 90, 150)
ORANGE = RGBColor(218, 112, 28)


def set_text(shape, text: str, size=16, bold=False, color=BLACK, font="Apple SD Gothic Neo"):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return p


def add_textbox(slide, x, y, w, h, text, size=16, bold=False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = set_text(box, text, size=size, bold=bold, color=color)
    if align:
        p.alignment = align
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.45, 0.25, 12.4, 0.45, title, size=25, bold=True)
    if subtitle:
        add_textbox(slide, 0.48, 0.72, 12.0, 0.30, subtitle, size=10.5, color=MID)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.05), Inches(12.4), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.color.rgb = BLACK


def add_card(slide, x, y, w, h, title, body, color=BLACK, title_size=13, body_size=10.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(220, 220, 220)
    add_textbox(slide, x + 0.15, y + 0.12, w - 0.3, 0.25, title, size=title_size, bold=True, color=color)
    tb = add_textbox(slide, x + 0.15, y + 0.48, w - 0.3, h - 0.6, body, size=body_size)
    tf = tb.text_frame
    tf.word_wrap = True
    return shape


def add_bullets(slide, x, y, w, h, items, size=13, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Apple SD Gothic Neo"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_badge(slide, x, y, text, color=BLACK, w=1.35):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.color.rgb = color
    add_textbox(slide, x, y + 0.07, w, 0.18, text, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return badge


def one_page_summary(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLACK
    bar.line.color.rgb = BLACK
    add_textbox(slide, 0.45, 0.18, 8.5, 0.35, "adiFit: AI Review Intelligence for adidas eCommerce", size=20, bold=True, color=WHITE)
    add_textbox(slide, 10.3, 0.21, 2.5, 0.25, "Name: Hyeoky Yun", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    add_card(
        slide,
        0.45,
        1.0,
        4.05,
        1.55,
        "1. WHAT YOU BUILT",
        "고객 리뷰에서 사이즈, 착용감, 색상/디자인, 사용 용도, VOC 이슈를 구조화하는 AI 리뷰 분석 도구입니다. 소비자에게는 근거 기반 핏 조언을, 머천다이저에게는 PDP 개선 액션을 제공합니다.",
        BLUE,
    )
    add_card(
        slide,
        4.65,
        1.0,
        4.05,
        1.55,
        "2. HOW IT WORKS",
        f"Demo reviews → Gemini structured extraction → embeddings/RAG → Streamlit dashboard. 추가로 KMeans/PCA로 고객 사용 맥락 세그먼트를 분석합니다.\nRepo: {REPO_URL}",
        GREEN,
        body_size=9.5,
    )
    add_card(
        slide,
        8.85,
        1.0,
        4.05,
        1.55,
        "3. KEY RESULT",
        "24개 데모 리뷰, 3개 상품 분석. 사이즈 신호 정확도 95.8%, 테스트 질문 citation coverage 100%. Samba OG는 발볼/사이즈업 리스크 세그먼트로 분리되었습니다.",
        ORANGE,
    )
    add_card(
        slide,
        0.45,
        2.85,
        6.1,
        1.55,
        "4. KEY CHALLENGES & HOW YOU SOLVED THEM",
        "LLM 환각과 근거 없는 사이즈 추천을 막기 위해 검색된 리뷰 안에서만 답하도록 제한하고 review_id 인용을 강제했습니다. 소량 데이터 한계는 직접 라벨링 gold set과 샘플 기반 검증으로 투명하게 표시했습니다.",
        RED,
    )
    add_card(
        slide,
        6.8,
        2.85,
        6.1,
        1.55,
        "5. LIMITATIONS & NEXT STEPS",
        "현재는 직접 정의한 24개 데모 리뷰 기반입니다. 실제 운영에서는 대규모 first-party 리뷰, 반품/CS 데이터, 상품 카탈로그를 연결해 세그먼트 drift와 VOC spike를 추적할 수 있습니다.",
        BLACK,
    )

    add_textbox(slide, 0.55, 4.75, 12.0, 0.28, "Core demo flow: Fit Advisor → Merchandiser Dashboard → Product Comparison → Segment & Trend → Validation", size=12, bold=True)
    return slide


def build_deck(path: Path, only_summary=False):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    one_page_summary(prs)

    if only_summary:
        prs.save(path)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Recorded Demo Video", "3분 이내 영상은 이 슬라이드에 삽입하거나 별도 파일로 첨부")
    add_card(
        slide,
        0.7,
        1.35,
        5.8,
        4.7,
        "Demo Flow",
        "0:00 문제와 솔루션 소개\n0:20 앱 실행 및 데이터 범위\n0:40 핏 어드바이저 질문 입력\n1:20 머천다이저 VOC 대시보드\n1:55 상품 비교\n2:15 세그먼트 & VOC 트렌드\n2:45 검증 탭\n3:00 마무리",
        BLUE,
        body_size=15,
    )
    add_card(
        slide,
        6.85,
        1.35,
        5.8,
        4.7,
        "Video Placeholder",
        "PPT 제출 전 여기에 녹화 영상을 삽입하거나, adiFit_demo.mp4를 별도 첨부하세요.\n\n추천 데모 질문:\n\"평소 260 신고 발볼 넓은데 정사이즈로 사도 될까요?\"\n상품: Samba OG",
        BLACK,
        body_size=15,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem & Topic Fit", "Assignment Topic 1: Consumer Review Analysis")
    add_bullets(
        slide,
        0.7,
        1.35,
        6.0,
        4.8,
        [
            "상품 상세 페이지에는 리뷰가 많지만 사이즈, 착용감, 색상, 용도, 품질 이슈가 섞여 있어 소비자와 내부 담당자가 빠르게 이해하기 어렵습니다.",
            "과제 요구사항의 고객 리뷰 분석을 선택하고, adidas footwear 도메인에 맞춰 fit/usage/VOC intelligence로 구체화했습니다.",
            "핵심 기능은 실제 동작합니다. UI는 Streamlit으로 구현했고, LLM 추출·검색·집계·클러스터링 결과를 바로 시연할 수 있습니다.",
        ],
        size=16,
    )
    add_card(slide, 7.2, 1.4, 5.3, 1.1, "Consumer Pain", "정사이즈인지, 발볼이 맞는지, 어떤 용도에 적합한지 리뷰를 일일이 찾아야 함", RED, body_size=12)
    add_card(slide, 7.2, 2.8, 5.3, 1.1, "Business Pain", "핏 불만, 내구성, 배송/포장 이슈를 수동으로 읽고 분류해야 함", ORANGE, body_size=12)
    add_card(slide, 7.2, 4.2, 5.3, 1.1, "AI Opportunity", "리뷰를 구조화해 PDP 개선, 추천 메시지, CS/VOC 대응으로 연결", GREEN, body_size=12)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Solution Architecture", "One analysis engine, two user-facing outputs")
    steps = [
        ("1", "Demo Reviews", "24 curated adidas-style reviews\n3 products"),
        ("2", "LLM Extraction", "size, width, aspects,\nuse case, comfort, VOC tags"),
        ("3", "Embedding/RAG", "Gemini embeddings\ncosine top-k search"),
        ("4", "Outputs", "fit advisor, dashboard,\nsegments, trend, evaluation"),
    ]
    x = 0.65
    for num, title, body in steps:
        add_badge(slide, x, 1.45, num, BLACK, w=0.45)
        add_card(slide, x, 1.9, 2.8, 1.65, title, body, BLUE if num in ("2", "3") else BLACK, body_size=12)
        x += 3.1
    add_bullets(
        slide,
        0.8,
        4.25,
        12.0,
        1.5,
        [
            "Core AI: Gemini structured output for review extraction and Gemini embeddings for retrieval.",
            "Engineering contract: Pydantic schemas keep extraction, aggregation, advisor, and UI consistent.",
            "Research layer: KMeans/PCA clustering and monthly tag aggregation expose usage segments and VOC drift.",
        ],
        size=14,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Structured Review Extraction", "Reviews are converted into machine-readable product intelligence")
    add_card(slide, 0.65, 1.35, 4.0, 4.9, "Signals", "size_signal\nwidth_signal\naspect sentiment\nuse_cases\ncomfort_tags\ndesign_color_tags\ncustomer_tags\nvoc_issue_tags", BLUE, body_size=15)
    add_card(slide, 4.95, 1.35, 3.8, 4.9, "Example", "Review: \"발볼 넓은 편이라 걱정했는데 255에서 260으로 한 치수 올리니 편해요.\"\n\nExtracted:\n- size: runs_small\n- width: wide\n- issue: fit_risk\n- customer: wide_feet", GREEN, body_size=12.5)
    add_card(slide, 9.05, 1.35, 3.55, 4.9, "Why It Matters", "단순 요약을 넘어 리뷰를 집계, 검색, 세그먼트 분석, 시계열 추적에 사용할 수 있는 구조화 데이터로 바꿉니다.", ORANGE, body_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Consumer Fit Advisor", "Grounded answer generation with cited review evidence")
    add_card(slide, 0.65, 1.35, 5.8, 4.75, "Demo Question", "상품: Samba OG\n질문: \"평소 260 신고 발볼 넓은데 정사이즈로 사도 될까요?\"\n\nExpected answer:\n한 사이즈 업을 권장. 발볼 넓은 고객 리뷰에서 정사이즈 착용 시 조임/교환 이슈가 반복됨.", BLUE, body_size=14)
    add_card(slide, 6.8, 1.35, 5.8, 4.75, "Guardrails", "검색된 리뷰 근거 안에서만 답변\nreview_id와 원문 인용을 citations로 표시\n근거 부족/상충 시 confidence를 낮추고 caveat 표시\ncitation coverage를 실시간 계산", GREEN, body_size=14)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Merchandiser Dashboard", "VOC is translated into PDP and operations actions")
    add_bullets(
        slide,
        0.7,
        1.35,
        6.0,
        4.8,
        [
            "상품별 리뷰 수, 평균 평점, 사이즈/발볼 결론을 한눈에 확인합니다.",
            "속성별 감성 차트와 사이즈 신호 분포로 고객 불만이 어디에 몰리는지 봅니다.",
            "VOC → Action 규칙으로 PDP 문구, 대체 모델 추천, 포장/배송 점검 같은 실행 제안을 생성합니다.",
        ],
        size=15,
    )
    add_card(slide, 7.1, 1.35, 5.5, 1.2, "Example Action", "Samba OG: PDP 사이즈 영역에 '정사이즈보다 반 치수 크게' 안내 노출", RED, body_size=12)
    add_card(slide, 7.1, 2.85, 5.5, 1.2, "Expected Impact", "사이즈 미스매치 교환·반품 감소, 구매 전 의사결정 지원", GREEN, body_size=12)
    add_card(slide, 7.1, 4.35, 5.5, 1.2, "Download", "상품별 VOC 리포트를 Markdown으로 다운로드 가능", BLUE, body_size=12)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Research Strength Layer", "Clustering + time-series view for usage segments and VOC drift")
    add_card(slide, 0.65, 1.35, 3.8, 4.8, "Usage Segments", "Comfort-driven Buyers\nWide-foot Fit Risk\nPerformance Runners\nPost-purchase VOC Risk\nMixed Review Segment", BLUE, body_size=15)
    add_card(slide, 4.8, 1.35, 3.8, 4.8, "Clustering", "Review embeddings are clustered with KMeans and projected with PCA. LLM tags explain each cluster in human-readable business terms.", GREEN, body_size=13)
    add_card(slide, 8.95, 1.35, 3.8, 4.8, "Time-series", "Monthly tag aggregation tracks fit risk, durability, delivery/packaging, price resistance, and comfort complaints. At scale, this becomes VOC spike/drift monitoring.", ORANGE, body_size=13)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Validation", "Measured reliability, not just a polished demo")
    add_card(slide, 0.65, 1.35, 3.8, 4.75, "Extraction", "Gold set: 24 manually labeled reviews\nMetric: size_signal accuracy\nResult: 95.8%\nMismatch: 1 / 24", GREEN, body_size=15)
    add_card(slide, 4.85, 1.35, 3.8, 4.75, "Grounding", "Advisor responses cite review_ids.\nTested advisor cases showed citation coverage 100%.\nThis is grounding coverage, not full factuality.", BLUE, body_size=13)
    add_card(slide, 9.05, 1.35, 3.8, 4.75, "Demo Scope", "24 curated reviews are enough to prove the working pipeline. Statistical significance and drift detection require larger first-party review data.", ORANGE, body_size=13)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Limitations & Next Steps", "What is real now, and what scales next")
    add_bullets(
        slide,
        0.75,
        1.35,
        12.0,
        4.9,
        [
            "Current data is manually curated demo data, not live adidas production data.",
            "Current evaluation focuses on size signal accuracy and citation coverage; full factuality and business impact need larger test sets.",
            "Next: connect first-party reviews, return/CS/VOC data, product catalog, inventory, and real PDP event data.",
            "Next: monitor segment drift and VOC spike over time, then trigger PDP copy changes or CS alerts.",
            "Next: multilingual support for global eCommerce reviews.",
        ],
        size=16,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Portfolio / Relevant Capability", "Why this implementation reflects the role requirements")
    add_card(slide, 0.65, 1.35, 3.8, 4.7, "AI Engineering", "LLM structured output\nRAG with grounding\nPydantic contracts\nStreamlit demo app\nEvaluation layer", BLUE, body_size=15)
    add_card(slide, 4.8, 1.35, 3.8, 4.7, "Research Fit", "Time-series analysis\nClustering/segmentation\nPattern detection\nVOC trend monitoring\nData-driven interpretation", GREEN, body_size=15)
    add_card(slide, 8.95, 1.35, 3.8, 4.7, "Business Translation", "PDP action recommendations\nReturn-risk framing\nConsumer decision support\nMerchandiser dashboard\nClear limitations and next steps", ORANGE, body_size=15)

    prs.save(path)


def main():
    OUT_DIR.mkdir(exist_ok=True)
    build_deck(OUT_DIR / "adiFit_one_page_summary.pptx", only_summary=True)
    build_deck(OUT_DIR / "adiFit_assignment_deck.pptx", only_summary=False)


if __name__ == "__main__":
    main()
